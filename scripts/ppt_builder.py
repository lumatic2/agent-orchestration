"""
ppt_builder.py — Claude Code용 PPTX 생성 유틸리티

사용법:
    python3 ~/.claude/ppt_builder.py --input slides.json --output result.pptx

slides.json 형식:
[
  {
    "layout": "title",          // title | section | content | two_col | blank
    "title": "슬라이드 제목",
    "subtitle": "부제목",       // layout=title 전용
    "body": "본문 내용\n줄바꿈 지원",
    "left": "왼쪽 컬럼",       // layout=two_col 전용
    "right": "오른쪽 컬럼"     // layout=two_col 전용
  }
]
"""

import json
import argparse
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 브랜드 컬러 (커스텀 가능)
COLOR_PRIMARY   = RGBColor(0x1A, 0x1A, 0x2E)   # 진남색
COLOR_ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # 남색
COLOR_HIGHLIGHT = RGBColor(0x0F, 0x3A, 0x6E)   # 파란계열
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY      = RGBColor(0x4A, 0x4A, 0x4A)
COLOR_LIGHT     = RGBColor(0xF5, 0xF5, 0xF5)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=COLOR_WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, COLOR_PRIMARY)
    # 왼쪽 강조 바
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_HIGHLIGHT
    bar.line.fill.background()
    # 제목
    add_textbox(slide, data.get("title", ""), Inches(0.5), Inches(2.5),
                Inches(12), Inches(1.5), font_size=40, bold=True, align=PP_ALIGN.LEFT)
    # 부제목
    if data.get("subtitle"):
        add_textbox(slide, data["subtitle"], Inches(0.5), Inches(4.2),
                    Inches(10), Inches(0.8), font_size=20, color=RGBColor(0xCC, 0xCC, 0xCC))


def build_section_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_ACCENT)
    add_textbox(slide, data.get("title", ""), Inches(1), Inches(3),
                Inches(11), Inches(1.2), font_size=32, bold=True, align=PP_ALIGN.CENTER)


def build_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    # 상단 제목 바
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    add_textbox(slide, data.get("title", ""), Inches(0.4), Inches(0.2),
                Inches(12), Inches(0.8), font_size=22, bold=True)
    # 본문
    body = data.get("body", "")
    add_textbox(slide, body, Inches(0.6), Inches(1.4),
                Inches(12.1), Inches(5.8), font_size=16, color=COLOR_GRAY, wrap=True)


def build_two_col_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_WHITE)
    # 상단 제목 바
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    add_textbox(slide, data.get("title", ""), Inches(0.4), Inches(0.2),
                Inches(12), Inches(0.8), font_size=22, bold=True)
    # 구분선
    div = slide.shapes.add_shape(1, Inches(6.5), Inches(1.4), Inches(0.05), Inches(5.8))
    div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    div.line.fill.background()
    # 왼쪽
    add_textbox(slide, data.get("left", ""), Inches(0.4), Inches(1.5),
                Inches(5.8), Inches(5.7), font_size=15, color=COLOR_GRAY, wrap=True)
    # 오른쪽
    add_textbox(slide, data.get("right", ""), Inches(6.7), Inches(1.5),
                Inches(6.3), Inches(5.7), font_size=15, color=COLOR_GRAY, wrap=True)


BUILDERS = {
    "title":   build_title_slide,
    "section": build_section_slide,
    "content": build_content_slide,
    "two_col": build_two_col_slide,
}


def build_pptx(slides_data: list, output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_data in enumerate(slides_data):
        layout = slide_data.get("layout", "content")
        builder = BUILDERS.get(layout, build_content_slide)
        builder(prs, slide_data)
        print(f"  슬라이드 {i+1}: [{layout}] {slide_data.get('title', '')[:40]}")

    prs.save(output_path)
    print(f"\n저장 완료: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="PPTX 생성기")
    parser.add_argument("--input",  "-i", required=True, help="slides.json 경로")
    parser.add_argument("--output", "-o", required=True, help="출력 .pptx 경로")
    args = parser.parse_args()

    with open(args.input, encoding="utf-8") as f:
        slides_data = json.load(f)

    print(f"슬라이드 {len(slides_data)}장 생성 중...")
    build_pptx(slides_data, args.output)


if __name__ == "__main__":
    main()
